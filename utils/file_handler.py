import os
import io
import csv
import json
import telebot
import pdfplumber
from openpyxl import load_workbook
from telebot.types import Message
from utils.db import is_premium, log_file_analysis
from deep_translator import GoogleTranslator
from docx import Document
from bs4 import BeautifulSoup
from pptx import Presentation

MAX_FREE_SIZE_MB = 2
MAX_PREMIUM_SIZE_MB = 20

def register_file_handler(bot: telebot.TeleBot):

    @bot.message_handler(content_types=['document'])
    def handle_file(msg: Message):
        user_id = msg.from_user.id
        file_info = bot.get_file(msg.document.file_id)
        file_size = msg.document.file_size / (1024 * 1024)  # MB
        is_user_premium = is_premium(user_id)

        size_limit = MAX_PREMIUM_SIZE_MB if is_user_premium else MAX_FREE_SIZE_MB
        if file_size > size_limit:
            bot.reply_to(msg, f"❌ الحد الأقصى لحجم الملف هو {size_limit}MB")
            return

        try:
            downloaded_file = bot.download_file(file_info.file_path)
            file_stream = io.BytesIO(downloaded_file)
            filename = msg.document.file_name.lower()

            if filename.endswith(".txt"):
                content = file_stream.read().decode("utf-8")
            elif filename.endswith(".csv"):
                content = ""
                reader = csv.reader(io.StringIO(file_stream.read().decode("utf-8")))
                for row in reader:
                    content += ", ".join(row) + "\n"
            elif filename.endswith(".xlsx"):
                wb = load_workbook(file_stream, read_only=True)
                content = ""
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        row_data = [str(cell) if cell is not None else "" for cell in row]
                        content += ", ".join(row_data) + "\n"
            elif filename.endswith(".pdf"):
                content = ""
                with pdfplumber.open(file_stream) as pdf:
                    for page in pdf.pages:
                        content += page.extract_text() + "\n"
            elif filename.endswith(".docx"):
                doc = Document(file_stream)
                content = "\n".join([para.text for para in doc.paragraphs])
            elif filename.endswith(".html") or filename.endswith(".htm"):
                soup = BeautifulSoup(file_stream.read(), "html.parser")
                content = soup.get_text()
            elif filename.endswith(".pptx"):
                prs = Presentation(file_stream)
                content = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
            else:
                bot.reply_to(msg, "❌ نوع الملف غير مدعوم.")
                return

            translated = GoogleTranslator(source='auto', target='ar').translate(content[:2000])
            bot.reply_to(msg, f"📝 ترجمة مقتطف من الملف:\n\n{translated}")
            log_file_analysis(user_id, filename, len(content))

        except Exception as e:
            bot.reply_to(msg, f"❌ حدث خطأ أثناء قراءة الملف:\n{str(e)}")
